import os
import io
import json
import tempfile
from pathlib import Path
from typing import List
from lxml import etree
from langchain_core.documents import Document
from langchain_community.document_loaders import TextLoader


# ── TXT ──────────────────────────────────────────────────────────────────────

def load_txt_file(file_path: str) -> List[Document]:
    try:
        loader = TextLoader(file_path, encoding="utf-8")
        return loader.load()
    except UnicodeDecodeError:
        loader = TextLoader(file_path, encoding="latin-1")
        return loader.load()
    except Exception as e:
        print(f"Error loading {file_path}: {e}")
        return []


# ── XML ──────────────────────────────────────────────────────────────────────

def load_xml_file(file_path: str) -> List[Document]:
    try:
        tree = etree.parse(file_path)
        root = tree.getroot()
        text_content = []
        for element in root.iter():
            if element.text and element.text.strip():
                text_content.append(element.text.strip())
            if element.tail and element.tail.strip():
                text_content.append(element.tail.strip())
        full_text = "\n".join(text_content)
        if full_text.strip():
            return [Document(page_content=full_text, metadata={"source": file_path, "type": "xml"})]
        return []
    except Exception as e:
        print(f"Error loading {file_path}: {e}")
        return []


# ── PDF ──────────────────────────────────────────────────────────────────────

def load_pdf_file(file_path: str) -> List[Document]:
    try:
        import pypdf
        reader = pypdf.PdfReader(file_path)
        docs = []
        for i, page in enumerate(reader.pages):
            text = page.extract_text() or ""
            if text.strip():
                docs.append(Document(
                    page_content=text,
                    metadata={"source": file_path, "type": "pdf", "page": i + 1}
                ))
        return docs
    except Exception as e:
        print(f"Error loading PDF {file_path}: {e}")
        return []


# ── CSV ──────────────────────────────────────────────────────────────────────

def load_csv_file(file_path: str) -> List[Document]:
    try:
        import pandas as pd
        df = pd.read_csv(file_path)
        text = df.to_string(index=False)
        return [Document(page_content=text, metadata={"source": file_path, "type": "csv"})]
    except Exception as e:
        print(f"Error loading CSV {file_path}: {e}")
        return []


# ── JSON ─────────────────────────────────────────────────────────────────────

def load_json_file(file_path: str) -> List[Document]:
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            data = json.load(f)
        text = json.dumps(data, indent=2, ensure_ascii=False)
        return [Document(page_content=text, metadata={"source": file_path, "type": "json"})]
    except Exception as e:
        print(f"Error loading JSON {file_path}: {e}")
        return []


# ── XLSX / XLS ───────────────────────────────────────────────────────────────

def load_excel_file(file_path: str) -> List[Document]:
    try:
        import pandas as pd
        xl = pd.ExcelFile(file_path)
        docs = []
        for sheet in xl.sheet_names:
            df = xl.parse(sheet)
            text = f"Sheet: {sheet}\n{df.to_string(index=False)}"
            if text.strip():
                docs.append(Document(
                    page_content=text,
                    metadata={"source": file_path, "type": "excel", "sheet": sheet}
                ))
        return docs
    except Exception as e:
        print(f"Error loading Excel {file_path}: {e}")
        return []


# ── DOCX ─────────────────────────────────────────────────────────────────────

def load_docx_file(file_path: str) -> List[Document]:
    try:
        from docx import Document as DocxDocument
        doc = DocxDocument(file_path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        full_text = "\n".join(paragraphs)
        if full_text.strip():
            return [Document(page_content=full_text, metadata={"source": file_path, "type": "docx"})]
        return []
    except Exception as e:
        print(f"Error loading DOCX {file_path}: {e}")
        return []


# ── Markdown ─────────────────────────────────────────────────────────────────

def load_md_file(file_path: str) -> List[Document]:
    return load_txt_file(file_path)  # plain-text read is fine for markdown


# ── HTML ─────────────────────────────────────────────────────────────────────

def load_html_file(file_path: str) -> List[Document]:
    try:
        from bs4 import BeautifulSoup
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            soup = BeautifulSoup(f.read(), "html.parser")
        text = soup.get_text(separator="\n")
        if text.strip():
            return [Document(page_content=text, metadata={"source": file_path, "type": "html"})]
        return []
    except Exception as e:
        print(f"Error loading HTML {file_path}: {e}")
        return []


# ── PPTX ─────────────────────────────────────────────────────────────────────

def load_pptx_file(file_path: str) -> List[Document]:
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        docs = []
        for i, slide in enumerate(prs.slides, 1):
            lines = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = para.text.strip()
                        if line:
                            lines.append(line)
            if lines:
                docs.append(Document(
                    page_content="\n".join(lines),
                    metadata={"source": file_path, "type": "pptx", "slide": i}
                ))
        return docs
    except Exception as e:
        print(f"Error loading PPTX {file_path}: {e}")
        return []


# ── DISPATCH MAP ─────────────────────────────────────────────────────────────

_EXT_LOADERS = {
    ".txt":  load_txt_file,
    ".md":   load_md_file,
    ".xml":  load_xml_file,
    ".pdf":  load_pdf_file,
    ".csv":  load_csv_file,
    ".json": load_json_file,
    ".xlsx": load_excel_file,
    ".xls":  load_excel_file,
    ".docx": load_docx_file,
    ".html": load_html_file,
    ".htm":  load_html_file,
    ".pptx": load_pptx_file,
}

SUPPORTED_EXTENSIONS = set(_EXT_LOADERS.keys())


def load_file_by_extension(file_path: str) -> List[Document]:
    """Dispatch to the correct loader based on file extension."""
    ext = Path(file_path).suffix.lower()
    loader_fn = _EXT_LOADERS.get(ext)
    if loader_fn:
        return loader_fn(file_path)
    print(f"Unsupported file type: {ext} ({file_path})")
    return []


# ── DIRECTORY LOADING ─────────────────────────────────────────────────────────

def load_all_documents(directory: Path) -> List[Document]:
    """Load all supported file types from a directory."""
    documents = []
    for ext in SUPPORTED_EXTENSIONS:
        for file_path in directory.glob(f"*{ext}"):
            docs = load_file_by_extension(str(file_path))
            documents.extend(docs)
    return documents


# ── UPLOADED FILE (in-memory) ─────────────────────────────────────────────────

def load_uploaded_file(file_content: bytes, filename: str) -> List[Document]:
    """Load a document from uploaded file bytes."""
    ext = Path(filename).suffix.lower()

    # Plain-text formats — decode and create Document directly (no temp file needed)
    if ext in (".txt", ".md"):
        try:
            text = file_content.decode("utf-8")
        except UnicodeDecodeError:
            text = file_content.decode("latin-1")
        return [Document(page_content=text, metadata={"source": filename, "type": ext.lstrip(".")})]

    if ext == ".xml":
        try:
            text = file_content.decode("utf-8")
        except UnicodeDecodeError:
            text = file_content.decode("latin-1")
        try:
            root = etree.fromstring(text.encode())
            text_content = []
            for element in root.iter():
                if element.text and element.text.strip():
                    text_content.append(element.text.strip())
                if element.tail and element.tail.strip():
                    text_content.append(element.tail.strip())
            full_text = "\n".join(text_content)
            if full_text.strip():
                return [Document(page_content=full_text, metadata={"source": filename, "type": "xml"})]
        except Exception as e:
            print(f"XML parse error for {filename}: {e}")
        return []

    if ext == ".json":
        try:
            text = file_content.decode("utf-8")
            data = json.loads(text)
            out = json.dumps(data, indent=2, ensure_ascii=False)
            return [Document(page_content=out, metadata={"source": filename, "type": "json"})]
        except Exception as e:
            print(f"JSON parse error for {filename}: {e}")
            return []

    if ext == ".html" or ext == ".htm":
        try:
            from bs4 import BeautifulSoup
            soup = BeautifulSoup(file_content, "html.parser")
            text = soup.get_text(separator="\n")
            if text.strip():
                return [Document(page_content=text, metadata={"source": filename, "type": "html"})]
        except Exception as e:
            print(f"HTML parse error for {filename}: {e}")
        return []

    # Binary formats — write to a temp file and reuse file-based loaders
    if ext not in (".pdf", ".csv", ".xlsx", ".xls", ".docx", ".pptx"):
        print(f"Unsupported file type for upload: {ext}")
        return []

    try:
        suffix = ext  # e.g. ".pdf"
        with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
            tmp.write(file_content)
            tmp_path = tmp.name

        docs = load_file_by_extension(tmp_path)
        # Replace temp path with original filename in metadata
        for doc in docs:
            doc.metadata["source"] = filename
        return docs
    except Exception as e:
        print(f"Error loading uploaded file {filename}: {e}")
        return []
    finally:
        try:
            os.unlink(tmp_path)
        except Exception:
            pass
